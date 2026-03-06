from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def inspect_presentation(pptx_path):
    prs = Presentation(pptx_path)

    for slide_index, slide in enumerate(prs.slides, start=1):
        print(f"\n--- Slide {slide_index} ---")

        for shape_index, shape in enumerate(slide.shapes, start=1):
            print(f"\nShape {shape_index}")
            print(f"  Type: {shape.shape_type}")
            print(f"  Name: {shape.name}")
            print(f"  Left: {shape.left}")
            print(f"  Top: {shape.top}")
            print(f"  Width: {shape.width}")
            print(f"  Height: {shape.height}")

            # TEXT
            if shape.has_text_frame:
                print("  Text:")
                for p_index, paragraph in enumerate(shape.text_frame.paragraphs):
                    print(f"    Paragraph {p_index + 1}: '{paragraph.text}'")
                    print(f"      Alignment: {paragraph.alignment}")

                    for r_index, run in enumerate(paragraph.runs):
                        font = run.font
                        print(f"        Run {r_index + 1}: '{run.text}'")
                        print(f"          Size: {font.size}")
                        print(f"          Bold: {font.bold}")
                        print(f"          Italic: {font.italic}")
                        print(f"          Name: {font.name}")
                        print(f"          Color: {font.color.rgb if font.color and font.color.rgb else None}")

            # IMAGE
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                print("  Picture:")
                print(f"    Image size: {shape.image.size}")

            # FILL
            if hasattr(shape, "fill"):
                fill = shape.fill
                print(f"  Fill type: {fill.type}")

            # LINE
            if hasattr(shape, "line"):
                line = shape.line
                print(f"  Line width: {line.width}")

if __name__ == "__main__":
    inspect_presentation("slide_data/background_template.pptx")
