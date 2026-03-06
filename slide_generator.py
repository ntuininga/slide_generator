from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import json
import os

BASE_DIR = os.path.dirname(__file__)
FORMATTING_FILE = os.path.join(BASE_DIR, "slide_data", "song_slide_formatting.json")
SONG_DIR = os.path.join(BASE_DIR, "song_data")
TEMPLATE_FILE = os.path.join(BASE_DIR, "slide_data", "intro_template.pptx")

content = {
    "HEADER": "Cornerstone United Reformed Church of Edmonton",
    "DATE": "February 15, 2026",
    "SERVICE": "Morning Service",
    "SCRIPTURE": "Genesis 48",
    "TITLE": "A wonderful title",
    "PASTOR": "Pastor Vander Lei"
}

# --- Slide order constants ---
INTRO_IDX    = 0
OFFERING_IDX = 1
CREED_IDX    = 2


def generate_powerpoint_slides(service_data):
    print(service_data)
    prs = Presentation(TEMPLATE_FILE)

    # Verify template has exactly 3 slides
    if len(prs.slides) != 3:
        print(f"[WARNING] Expected 3 template slides, got {len(prs.slides)}")

    songs = service_data.get("songs", [])
    if not songs:
        print("[WARNING] No songs found in service data.")

    include_creed = service_data.get("includeCreed", False)

    # If no creed slide, remove it from the template
    if not include_creed:
        _remove_slide(prs, CREED_IDX)

    # Split songs: all but the last go before offerings, last goes after creed (or offerings)
    songs_before = songs[:-1] if len(songs) > 1 else songs
    last_song = songs[-1] if len(songs) > 1 else None

    # --- 1. Replace intro slide ---
    service_content = content.copy()
    service_content["DATE"]      = service_data.get("date", "")
    service_content["TITLE"]     = service_data.get("title", "")
    service_content["SERVICE"]   = service_data.get("service_label", "")
    service_content["SCRIPTURE"] = service_data.get("scripture", "")
    replace_slide_placeholders(prs.slides[INTRO_IDX], service_content)

    # --- 2. Replace offering slide ---
    replace_slide_placeholders(prs.slides[OFFERING_IDX], {
        "OFFERING": service_data.get("offering", "")
    })

    # --- 3. Insert songs before offerings slide (after intro, at index 1) ---
    insert_position = 1  # right after intro
    for song in reversed(songs_before):
        song_number = song.get("number")
        song_path = os.path.join(SONG_DIR, f"song_{song_number}.json")
        if not os.path.exists(song_path):
            print(f"[WARNING] Song #{song_number} does not exist. Skipping.")
            continue
        insert_song_slides(prs, insert_position, song_path, song.get("verses"), FORMATTING_FILE)

    # --- 4. Append last song after creed/offerings ---
    if last_song:
        song_number = last_song.get("number")
        song_path = os.path.join(SONG_DIR, f"song_{song_number}.json")
        if not os.path.exists(song_path):
            print(f"[WARNING] Last song #{song_number} does not exist. Skipping.")
        else:
            append_song_slides(prs, song_path, last_song.get("verses"), FORMATTING_FILE)

    prs.save(f"{service_data['service_label']}.pptx")
    print(f"[INFO] Saved: {service_data['service_label']}.pptx")


def _remove_slide(prs, index):
    """Remove a slide by index from the presentation."""
    xml_slides = prs.slides._sldIdLst
    slide_elem = xml_slides[index]
    xml_slides.remove(slide_elem)


def insert_song_slides(prs, position, json_file, verses, formatting_file):
    """Generate song slides and insert them at a specific index."""
    new_slides = _build_song_slides(prs, json_file, verses, formatting_file)
    # Insert in reverse so first verse slide ends up at `position`
    for slide in reversed(new_slides):
        _move_slide(prs, slide, position)


def append_song_slides(prs, json_file, verses, formatting_file):
    """Generate song slides and append them at the end."""
    _build_song_slides(prs, json_file, verses, formatting_file)
    # _build_song_slides already appends to end by default


def _build_song_slides(prs, json_file, verses, formatting_file):
    """Create all slides for a song, appended to the end. Returns list of new slides."""
    with open(json_file, "r", encoding="utf-8") as f:
        song = json.load(f)
    with open(formatting_file, "r", encoding="utf-8") as f:
        fmt = json.load(f)

    copyright_text = " | ".join(song["copyright"]) if song.get("copyright") else None

    # Resolve verses
    if verses is None:
        resolved_verses = song["verses"]
    else:
        verse_nums = [int(v) for v in verses]
        resolved_verses = [v for v in song["verses"] if v["verse"] in verse_nums]

    header_width    = Inches(float(fmt["header_width"]))
    header_font_size = Pt(float(fmt["header_font_size"]))
    verse_width     = Inches(float(fmt["verse_width"]))
    verse_font_size  = Pt(float(fmt["verse_font_size"]))

    slide_layout = prs.slide_layouts[6]
    slide_width  = prs.slide_width
    verse_left   = (slide_width - verse_width) / 2
    header_left  = (slide_width - header_width) / 2

    created_slides = []

    # Title slide (first verse)
    slide = prs.slides.add_slide(slide_layout)
    _clear_placeholders(slide)
    created_slides.append(slide)

    _add_textbox(slide, header_left, Inches(0.5), header_width, Inches(2),
                 f"#{song['number']} {song['title']}",
                 header_font_size, bold=True, align=PP_ALIGN.CENTER)

    first_verse = resolved_verses[0] if resolved_verses else song["verses"][0]
    _add_textbox(slide, verse_left, Inches(2.5), verse_width, Inches(4),
                 f"{first_verse['verse']}. {first_verse['text']}",
                 verse_font_size, align=PP_ALIGN.CENTER)

    if copyright_text:
        _add_copyright(slide, prs, copyright_text)

    # Remaining verses
    for verse in resolved_verses[1:]:
        slide = prs.slides.add_slide(slide_layout)
        _clear_placeholders(slide)
        created_slides.append(slide)

        tb = _add_textbox(slide, verse_left, Inches(1), verse_width, Inches(5),
                          f"{verse['verse']}. {verse['text']}",
                          verse_font_size, align=PP_ALIGN.CENTER)
        tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    return created_slides


def _move_slide(prs, slide, target_index):
    """Move a slide (already in the presentation) to target_index."""
    xml_slides = prs.slides._sldIdLst
    # Find the slide's current element
    slide_id = None
    for sldId in xml_slides:
        if prs.slides._sldIdLst.index(sldId) == prs.slides.index(slide):
            slide_id = sldId
            break

    # Fallback: get last added slide element
    slide_id = xml_slides[-1]
    xml_slides.remove(slide_id)
    xml_slides.insert(target_index, slide_id)


def _clear_placeholders(slide):
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            shape._element.getparent().remove(shape._element)


def _add_textbox(slide, left, top, width, height, text, font_size,
                 bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    par = tf.paragraphs[0]
    par.font.size = font_size
    par.font.bold = bold
    par.alignment = align
    return tb


def _add_copyright(slide, prs, copyright_text):
    cb = slide.shapes.add_textbox(
        Inches(0.5),
        prs.slide_height - Inches(0.7),
        prs.slide_width - Inches(1),
        Inches(0.5)
    )
    tf = cb.text_frame
    tf.word_wrap = True
    tf.text = copyright_text
    par = tf.paragraphs[0]
    par.font.size = Pt(12)
    par.font.italic = True
    par.font.color.rgb = RGBColor(120, 120, 120)
    par.alignment = PP_ALIGN.CENTER


def replace_slide_placeholders(slide, replacements):
    """Replace placeholder text tokens in a slide's text runs."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for key, value in replacements.items():
                    if value and key in run.text:
                        run.text = run.text.replace(key, value)