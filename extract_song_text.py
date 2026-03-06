import json
from pptx import Presentation
import re
import os

BASE_DIR = os.path.dirname(__file__)
SLIDES_DIR = os.path.join(BASE_DIR, "song_slides")
OUTPUT_DIR = os.path.join(BASE_DIR, "song_data")

def extract_text(ppt_path):
    prs = Presentation(ppt_path)
    slides_data = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_lines = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_lines.extend(shape.text.splitlines())

        slides_data.append({
            "slide_number": slide_number,
            "lines": slide_lines
        })

    return slides_data

def parse_song(slides):
    verses = []
    current_verse = None
    current_text = []
    song_number = None
    song_title = None
    copyright_lines = set()
    refrain_text = None  # default to None

    for slide in slides:
        for raw_line in slide["lines"]:
            line = raw_line.strip()
            if not line:
                continue

            # Detect header (#337 Title)
            header_match = re.match(r"#(\d+[A-Za-z]?)\s+(.+)", line)
            if header_match:
                song_number = header_match.group(1)
                song_title = header_match.group(2)
                continue

            # Extract inline copyright lines
            if "copyright" in line.lower() or "publications" in line.lower():
                copyright_lines.add(line)
                continue

            # Detect refrain (case-insensitive)
            refrain_match = re.match(r"refrain[:\s]*", line, re.IGNORECASE)
            if refrain_match:
                # Store refrain text after "refrain" keyword
                refrain_text = line[refrain_match.end():].strip()
                continue

            # Verse detection (allow 1-2 digit verse numbers)
            verse_match = re.match(r"^(\d{1,2})\.?\s+(.*)", line)
            if verse_match:
                verse_number = int(verse_match.group(1))

                # Ignore garbage large numbers (like 1956)
                if verse_number > 20:
                    continue

                # Save previous verse
                if current_verse is not None:
                    verses.append({
                        "verse": current_verse,
                        "text": " ".join(current_text).strip()
                    })

                current_verse = verse_number
                current_text = [verse_match.group(2)]
                continue

            # Continue accumulating verse text
            if current_verse is not None:
                current_text.append(line)

    # Save last verse
    if current_verse is not None:
        verses.append({
            "verse": current_verse,
            "text": " ".join(current_text).strip()
        })

    return {
        "number": song_number,
        "title": song_title,
        "verses": verses,
        "refrain": refrain_text,  # this will be None if no refrain
        "copyright": sorted(list(copyright_lines))
    }

def parse_folder(folder_path: str):
    master_data = {}
    for file_name in os.listdir(folder_path):
        if file_name.lower().endswith(".pptx"):
            text = extract_text(os.path.join(SLIDES_DIR, file_name))
            data = parse_song(text)
            master_data[data['number']] = data
            data_file_name = f"song_{data['number']}.json"
            save_to_json(data, os.path.join(OUTPUT_DIR, data_file_name))
        save_to_json(master_data, os.path.join(OUTPUT_DIR, "all_songs.json"))

def save_to_json(data, filename):
    with open(filename, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=4)

if __name__ == "__main__":
    parse_folder(SLIDES_DIR)
