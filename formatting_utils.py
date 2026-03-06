from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

def format_header(text_frame):
    text_frame.word_wrap = True 
    par = text_frame.paragraphs[0] 
    par.font.bold = True
    par.font.size = Pt(32)
    par.alignment = PP_ALIGN.CENTER

def format_body(text_frame):
    text_frame.word_wrap = True 
    par = text_frame.paragraphs[0] 
    par.font.size = Pt(25)
    par.alignment = PP_ALIGN.CENTER




